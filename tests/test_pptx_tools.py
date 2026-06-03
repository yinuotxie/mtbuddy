from __future__ import annotations

import json
import shutil
import subprocess
import sys
from pathlib import Path

from pptx import Presentation


REPO_ROOT = Path(__file__).resolve().parents[1]


def test_pptx_tool_creates_real_presentation_from_outline(tmp_path: Path) -> None:
    workspace = tmp_path / "demo"
    shutil.copytree(REPO_ROOT / "workspaces" / "demo", workspace)

    completed = subprocess.run(
        [sys.executable, "-m", "mtbuddy.mtclaw_tools", "mtbuddy_pptx_create_from_outline"],
        cwd=REPO_ROOT,
        input=json.dumps(
            {
                "workspace": str(workspace),
                "outline_path": "presentation-outline.md",
                "output_name": "competition-demo.pptx",
            }
        ),
        text=True,
        capture_output=True,
        check=False,
    )

    assert completed.returncode == 0, completed.stderr
    payload = json.loads(completed.stdout)
    assert payload["result"] == "ok"

    deck_path = Path(payload["presentation"])
    assert deck_path.exists()
    assert deck_path.suffix == ".pptx"

    presentation = Presentation(str(deck_path))
    slide_text = "\n".join(
        shape.text
        for slide in presentation.slides
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False)
    )
    assert len(presentation.slides) == 4
    assert "MTBUDDY Competition Demo" in slide_text
    assert "Why MTBUDDY" in slide_text
    assert "Capability Packs" in slide_text
    assert "Demo Flow" in slide_text


def test_pptx_tool_reads_text_from_generated_presentation(tmp_path: Path) -> None:
    workspace = tmp_path / "demo"
    shutil.copytree(REPO_ROOT / "workspaces" / "demo", workspace)

    create_completed = subprocess.run(
        [sys.executable, "-m", "mtbuddy.mtclaw_tools", "mtbuddy_pptx_create_from_outline"],
        cwd=REPO_ROOT,
        input=json.dumps(
            {
                "workspace": str(workspace),
                "outline_path": "presentation-outline.md",
            }
        ),
        text=True,
        capture_output=True,
        check=False,
    )
    assert create_completed.returncode == 0, create_completed.stderr

    read_completed = subprocess.run(
        [sys.executable, "-m", "mtbuddy.mtclaw_tools", "mtbuddy_pptx_read_text"],
        cwd=REPO_ROOT,
        input=json.dumps(
            {
                "workspace": str(workspace),
                "path": "artifacts/presentation.pptx",
            }
        ),
        text=True,
        capture_output=True,
        check=False,
    )

    assert read_completed.returncode == 0, read_completed.stderr
    payload = json.loads(read_completed.stdout)
    assert payload["result"] == "ok"
    assert len(payload["slides"]) == 4
    assert any(
        "OpenAI-compatible agents can call MTBUDDY" in "\n".join(slide["text"])
        for slide in payload["slides"]
    )
