"""PowerPoint generation and reading skill."""

from __future__ import annotations

from dataclasses import dataclass, field
from io import BytesIO
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR
from pptx.util import Inches, Pt

from ..core.interfaces import ArtifactStore, AuditLog
from ..core.workspace import ensure_inside_workspace


ALLOWED_OUTLINE_SUFFIXES = {".md", ".txt"}
DEFAULT_PPTX_NAME = "presentation.pptx"

PRIMARY = RGBColor(30, 39, 97)
SECONDARY = RGBColor(202, 220, 252)
ACCENT = RGBColor(2, 128, 144)
INK = RGBColor(33, 41, 92)
MUTED = RGBColor(91, 100, 128)
WHITE = RGBColor(255, 255, 255)
OFF_WHITE = RGBColor(247, 249, 252)


@dataclass(frozen=True)
class SlideOutline:
    """One generated slide worth of content."""

    title: str
    bullets: list[str] = field(default_factory=list)


@dataclass(frozen=True)
class DeckOutline:
    """Parsed Markdown outline for PPTX generation."""

    title: str
    slides: list[SlideOutline]


class PptxSkill:
    """Create and inspect simple, deterministic PowerPoint artifacts."""

    name = "pptx_skill"

    def __init__(
        self,
        workspace: Path,
        artifact_store: ArtifactStore,
        audit_log: AuditLog,
    ) -> None:
        self._workspace = workspace.resolve()
        self._artifact_store = artifact_store
        self._audit_log = audit_log

    def permissions(self) -> list[str]:
        return [
            f"Read .md, .txt, and .pptx files under {self._workspace}",
            "Write generated .pptx artifacts under the workspace artifacts directory.",
        ]

    def create_from_outline(
        self,
        outline_path: Path,
        *,
        output_name: str = DEFAULT_PPTX_NAME,
        title: str | None = None,
    ) -> Path:
        source_path = ensure_inside_workspace(self._workspace, outline_path)
        if source_path.suffix.lower() not in ALLOWED_OUTLINE_SUFFIXES:
            raise ValueError(f"outline file type is not allowed: {source_path.suffix}")

        outline = parse_markdown_outline(source_path.read_text(encoding="utf-8"), title=title)
        deck_bytes = render_presentation(outline)
        artifact_path = self._artifact_store.write_bytes(_pptx_name(output_name), deck_bytes)
        self._audit_log.record(
            "pptx.create_from_outline",
            details={
                "outline": source_path.relative_to(self._workspace).as_posix(),
                "path": str(artifact_path),
                "slides": len(outline.slides) + 1,
            },
        )
        return artifact_path

    def read_text(self, path: Path) -> dict[str, Any]:
        pptx_path = ensure_inside_workspace(self._workspace, path)
        if pptx_path.suffix.lower() != ".pptx":
            raise ValueError(f"presentation file type is not allowed: {pptx_path.suffix}")

        presentation = Presentation(str(pptx_path))
        slides: list[dict[str, Any]] = []
        for index, slide in enumerate(presentation.slides, start=1):
            texts = [
                shape.text.strip()
                for shape in slide.shapes
                if getattr(shape, "has_text_frame", False) and shape.text.strip()
            ]
            slides.append({"slide": index, "text": texts})

        self._audit_log.record(
            "pptx.read_text",
            details={
                "path": pptx_path.relative_to(self._workspace).as_posix(),
                "slides": len(slides),
            },
        )
        return {"result": "ok", "path": str(pptx_path), "slides": slides}


def parse_markdown_outline(content: str, *, title: str | None = None) -> DeckOutline:
    deck_title = title or "MTBUDDY Presentation"
    slides: list[SlideOutline] = []
    current_title: str | None = None
    current_bullets: list[str] = []

    for raw_line in content.splitlines():
        line = raw_line.strip()
        if not line:
            continue
        if line.startswith("# "):
            deck_title = title or line.removeprefix("# ").strip()
            continue
        if line.startswith("## "):
            if current_title:
                slides.append(SlideOutline(current_title, current_bullets))
            current_title = line.removeprefix("## ").strip()
            current_bullets = []
            continue
        if line.startswith(("- ", "* ")) and current_title:
            current_bullets.append(line[2:].strip())
            continue
        if current_title:
            current_bullets.append(line)

    if current_title:
        slides.append(SlideOutline(current_title, current_bullets))
    if not slides:
        slides.append(SlideOutline(deck_title, _fallback_bullets(content)))

    return DeckOutline(title=deck_title, slides=slides[:8])


def render_presentation(outline: DeckOutline) -> bytes:
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)

    _add_title_slide(presentation, outline.title, len(outline.slides))
    for index, slide_outline in enumerate(outline.slides, start=1):
        _add_content_slide(presentation, slide_outline, index, len(outline.slides))

    output = BytesIO()
    presentation.save(output)
    return output.getvalue()


def _add_title_slide(presentation: Presentation, title: str, content_slide_count: int) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _fill_background(slide, PRIMARY)
    _add_accent_band(slide, ACCENT, x=0, y=0, width=0.18, height=7.5)
    _add_circle(slide, SECONDARY, x=10.1, y=0.7, size=1.8)
    _add_circle(slide, ACCENT, x=11.1, y=5.35, size=0.75)
    _add_text(
        slide,
        title,
        x=0.85,
        y=1.65,
        width=8.8,
        height=1.4,
        size=40,
        color=WHITE,
        bold=True,
    )
    _add_text(
        slide,
        f"{content_slide_count} section deck generated by MTBUDDY",
        x=0.9,
        y=3.25,
        width=6.4,
        height=0.45,
        size=16,
        color=SECONDARY,
    )


def _add_content_slide(
    presentation: Presentation,
    outline: SlideOutline,
    index: int,
    total: int,
) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _fill_background(slide, OFF_WHITE)
    _add_accent_band(slide, PRIMARY, x=0, y=0, width=0.22, height=7.5)
    _add_circle(slide, ACCENT, x=0.65, y=0.55, size=0.62)
    _add_text(
        slide,
        f"{index:02d}",
        x=0.79,
        y=0.68,
        width=0.35,
        height=0.2,
        size=9,
        color=WHITE,
        bold=True,
    )
    _add_text(
        slide,
        outline.title,
        x=1.25,
        y=0.55,
        width=10.5,
        height=0.8,
        size=30,
        color=INK,
        bold=True,
    )

    bullets = outline.bullets[:5] or ["No bullet content was provided."]
    for bullet_index, bullet in enumerate(bullets):
        y = 1.65 + bullet_index * 0.82
        _add_bullet_card(slide, bullet, x=1.3, y=y, width=10.6, height=0.58)

    _add_text(
        slide,
        f"MTBUDDY artifact QA | slide {index} of {total}",
        x=1.25,
        y=6.92,
        width=5.5,
        height=0.25,
        size=10,
        color=MUTED,
    )


def _add_bullet_card(slide, text: str, *, x: float, y: float, width: float, height: float) -> None:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = SECONDARY
    _add_circle(slide, ACCENT, x=x + 0.18, y=y + 0.18, size=0.18)
    _add_text(
        slide,
        text,
        x=x + 0.55,
        y=y + 0.12,
        width=width - 0.8,
        height=height - 0.08,
        size=15,
        color=INK,
    )


def _fill_background(slide, color: RGBColor) -> None:
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_accent_band(slide, color: RGBColor, *, x: float, y: float, width: float, height: float) -> None:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _add_circle(slide, color: RGBColor, *, x: float, y: float, size: float) -> None:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(x),
        Inches(y),
        Inches(size),
        Inches(size),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _add_text(
    slide,
    text: str,
    *,
    x: float,
    y: float,
    width: float,
    height: float,
    size: int,
    color: RGBColor,
    bold: bool = False,
) -> None:
    textbox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(height))
    frame = textbox.text_frame
    frame.clear()
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    frame.vertical_anchor = MSO_ANCHOR.TOP
    paragraph = frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def _fallback_bullets(content: str) -> list[str]:
    lines = [line.strip("# -*\t ") for line in content.splitlines() if line.strip()]
    return lines[:5] or ["Create a concise presentation outline before generating slides."]


def _pptx_name(output_name: str) -> str:
    candidate = Path(output_name).name
    if not candidate.endswith(".pptx"):
        candidate = f"{candidate}.pptx"
    return candidate
